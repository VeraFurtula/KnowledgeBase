"""
RAG API: LangChain + Chroma + Ollama embeddings. Run from `server/`:
  pip install -r requirements.txt
  copy .env.example .env   # then edit
  python -m uvicorn app:app --reload --host 127.0.0.1 --port 8000
"""

from __future__ import annotations

import io
import os
import re
import threading
from pathlib import Path
from typing import Any

import chromadb
from dotenv import load_dotenv
from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from langchain_community.embeddings import OllamaEmbeddings
from langchain_community.vectorstores import Chroma
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pydantic import BaseModel, Field

try:
    from pptx import Presentation as _PptxPresentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE as _MSO_SHAPE_TYPE
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False

load_dotenv()


def _setup_langsmith() -> None:
    tracing = os.getenv("LANGSMITH_TRACING", "").lower() in ("true", "1", "yes")
    if tracing:
        os.environ["LANGCHAIN_TRACING_V2"] = "true"
    key = os.getenv("LANGSMITH_API_KEY", "").strip()
    if key:
        os.environ["LANGCHAIN_API_KEY"] = key
    proj = os.getenv("LANGSMITH_PROJECT", "").strip()
    if proj:
        os.environ["LANGCHAIN_PROJECT"] = proj
    endpoint = os.getenv("LANGSMITH_ENDPOINT", "").strip()
    if endpoint:
        os.environ["LANGCHAIN_ENDPOINT"] = endpoint


_setup_langsmith()

CHROMA_PATH = os.path.abspath(os.getenv("CHROMA_PATH", os.path.join(os.path.dirname(__file__), "data", "chroma")))

# Prevent concurrent /index calls from racing and corrupting the Chroma collection.
_index_lock = threading.Lock()
IMAGES_PATH = os.path.abspath(os.getenv("IMAGES_PATH", os.path.join(os.path.dirname(__file__), "data", "images")))
VECTOR_BACKEND = os.getenv("VECTOR_BACKEND", "chroma").lower().strip()
OLLAMA_BASE_URL = os.getenv("OLLAMA_BASE_URL", "http://127.0.0.1:11434").rstrip("/")
EMB_MODEL = os.getenv("EMB_MODEL", "nomic-embed-text").strip()

os.makedirs(CHROMA_PATH, exist_ok=True)
os.makedirs(IMAGES_PATH, exist_ok=True)

embeddings = OllamaEmbeddings(model=EMB_MODEL, base_url=OLLAMA_BASE_URL)


def _chunk_settings() -> tuple[int, int]:
    """eFront-style manuals: ~800–1600 chars window, 150–300 overlap (env-tunable).
    Larger default (1200) means each chunk carries more context per retrieval slot,
    so the model reads more document content within the same MAX_RETRIEVAL_CHARS budget."""
    size = int(os.getenv("CHUNK_SIZE", "1200"))
    overlap = int(os.getenv("CHUNK_OVERLAP", "250"))
    size = max(800, min(size, 1600))
    overlap = max(150, min(overlap, 300))
    if overlap >= size:
        overlap = max(150, min(250, size // 2))
    return size, overlap


CHUNK_SIZE, CHUNK_OVERLAP = _chunk_settings()
splitter = RecursiveCharacterTextSplitter(chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP)

# Markers emitted by the browser extractor (see `extractDocumentText.ts`).
_MARK_PAGE_OR_SLIDE = re.compile(r"(?:^|\n)--- (?:PDF page|PPTX slide) (\d+) ---\n")


def collection_name(user_id: str) -> str:
    safe = re.sub(r"[^a-zA-Z0-9_-]", "_", user_id.strip().lower())
    return (safe[:63] or "user") + "_kb"


def _safe_doc_name(filename: str) -> str:
    stem = Path(filename).stem
    return re.sub(r"[^a-zA-Z0-9_-]", "_", stem)[:60] or "doc"


def _extract_pptx_images(file_bytes: bytes, user_id: str, filename: str) -> list[dict[str, Any]]:
    """Extract embedded images from each PPTX slide and save them to disk."""
    if not _PPTX_AVAILABLE:
        return []
    safe = _safe_doc_name(filename)
    img_dir = os.path.join(IMAGES_PATH, user_id, safe)
    os.makedirs(img_dir, exist_ok=True)
    refs: list[dict[str, Any]] = []
    try:
        prs = _PptxPresentation(io.BytesIO(file_bytes))
        for slide_idx, slide in enumerate(prs.slides):
            slide_num = slide_idx + 1
            img_idx = 0
            for shape in slide.shapes:
                try:
                    if shape.shape_type != _MSO_SHAPE_TYPE.PICTURE:
                        continue
                    image = shape.image
                    ext = (image.ext or "png").lower()
                    if ext not in ("png", "jpg", "jpeg", "gif", "bmp", "emf", "wmf"):
                        ext = "png"
                    img_filename = f"slide_{slide_num}_{img_idx}.{ext}"
                    img_path = os.path.join(img_dir, img_filename)
                    with open(img_path, "wb") as f:
                        f.write(image.blob)
                    refs.append({
                        "source": filename,
                        "slide": slide_num,
                        "url": f"/images/{user_id}/{safe}/{img_filename}",
                    })
                    img_idx += 1
                except Exception:
                    pass
    except Exception:
        pass
    return refs


app = FastAPI(title="Knowledge Base RAG")
_cors = os.getenv("CORS_ORIGINS", "").strip()
app.add_middleware(
    CORSMiddleware,
    allow_origins=_cors.split(",") if _cors else ["*"],
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
)


class IndexDoc(BaseModel):
    id: str
    text: str
    filename: str = ""


class IndexRequest(BaseModel):
    user_id: str = Field(..., min_length=1, max_length=256)
    documents: list[IndexDoc]


class SearchRequest(BaseModel):
    user_id: str = Field(..., min_length=1, max_length=256)
    query: str = Field(..., min_length=1, max_length=32_000)
    k: int = Field(12, ge=1, le=40)


def _delete_collection(name: str) -> None:
    client = chromadb.PersistentClient(path=CHROMA_PATH)
    try:
        client.delete_collection(name)
    except Exception:
        pass


def _marked_blocks(text: str) -> list[tuple[int | None, str]]:
    """Split on PDF page / PPTX slide markers so chunk metadata can carry page/slide numbers."""
    matches = list(_MARK_PAGE_OR_SLIDE.finditer(text))
    if not matches:
        return [(None, text)]
    blocks: list[tuple[int | None, str]] = []
    first = matches[0].start()
    if first > 0:
        head = text[:first].strip()
        if head:
            blocks.append((None, head))
    for i, m in enumerate(matches):
        page = int(m.group(1))
        a = m.end()
        b = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        body = text[a:b].strip()
        if body:
            blocks.append((page, body))
    return blocks if blocks else [(None, text)]


def _section_hint(chunk: str) -> str:
    """First substantive line as a weak 'section' label for citations."""
    for line in chunk.splitlines():
        t = line.strip()
        if len(t) < 4:
            continue
        return t[:140] + ("…" if len(t) > 140 else "")
    return ""


def _documents_to_langchain_docs(rows: list[IndexDoc]) -> list[Document]:
    out: list[Document] = []
    g = 0
    for row in rows:
        text = (row.text or "").strip()
        if not text:
            continue
        src = row.filename or row.id
        module = (Path(src).stem[:80] or "document")
        for page_num, block in _marked_blocks(text):
            for chunk in splitter.split_text(block):
                chunk = chunk.strip()
                if not chunk:
                    continue
                meta: dict[str, Any] = {
                    "source": src,
                    "doc_id": row.id,
                    "chunk": g,
                    "module": module,
                    "section": _section_hint(chunk),
                }
                if page_num is not None:
                    meta["page"] = page_num
                out.append(Document(page_content=chunk, metadata=meta))
                g += 1
    return out


@app.get("/")
def root() -> dict[str, Any]:
    """Browser default is GET `/` — avoid a bare 404 when developers open the RAG port."""
    return {
        "service": "Knowledge Base RAG API",
        "note": "There is no page at `/`. Use the JSON endpoints below or open `/docs` for interactive API.",
        "endpoints": {
            "GET /health": "status, embedding model, chunk settings",
            "POST /index": "body: { user_id, documents: [{ id, text, filename }] }",
            "POST /search": "body: { user_id, query, k }",
            "GET /docs": "Swagger UI (try /health here)",
        },
    }


@app.post("/extract-images")
async def extract_images(
    user_id: str = Form(...),
    filename: str = Form(...),
    file: UploadFile = File(...),
) -> dict[str, Any]:
    """Extract and store images from an uploaded PPTX file."""
    uid = re.sub(r"[^a-zA-Z0-9_@.-]", "_", user_id.strip().lower())[:120]
    content = await file.read()
    refs = _extract_pptx_images(content, uid, filename)
    return {"image_refs": refs, "count": len(refs), "pptx_available": _PPTX_AVAILABLE}


@app.get("/images/{user_id}/{safe_name}/{img_name}")
def serve_image(user_id: str, safe_name: str, img_name: str) -> FileResponse:
    """Serve a previously extracted slide image."""
    path = os.path.join(IMAGES_PATH, user_id, safe_name, img_name)
    if not os.path.isfile(path):
        raise HTTPException(404, "Image not found")
    return FileResponse(path)


@app.get("/health")
def health() -> dict[str, Any]:
    return {
        "ok": True,
        "vector_backend": VECTOR_BACKEND,
        "chroma_path": CHROMA_PATH,
        "emb_model": EMB_MODEL,
        "ollama_base": OLLAMA_BASE_URL,
        "chunk_size": CHUNK_SIZE,
        "chunk_overlap": CHUNK_OVERLAP,
    }


@app.post("/index")
def index_docs(body: IndexRequest) -> dict[str, str]:
    if VECTOR_BACKEND != "chroma":
        raise HTTPException(400, f"Unsupported VECTOR_BACKEND={VECTOR_BACKEND!r} (only 'chroma').")

    name = collection_name(body.user_id)
    # Serialize all index operations for this process — prevents concurrent calls from
    # racing (delete + rebuild) and leaving the collection with fewer docs than intended.
    with _index_lock:
        try:
            docs = _documents_to_langchain_docs(body.documents)
            _delete_collection(name)

            if not docs:
                return {"status": "empty", "collection": name}

            Chroma.from_documents(
                documents=docs,
                embedding=embeddings,
                persist_directory=CHROMA_PATH,
                collection_name=name,
            )
        except Exception as exc:
            import traceback
            tb = traceback.format_exc()
            print(f"INDEX ERROR: {exc}\n{tb}", flush=True)
            raise HTTPException(500, detail=f"Index failed: {type(exc).__name__}: {exc}")
    return {"status": "indexed", "collection": name, "chunks": str(len(docs))}


def _multi_source_hits(raw_coll: Any, query_embedding: list, sources: list[str], per_src_k: int) -> list[Document]:
    """Query each source document separately using a pre-computed embedding.
    Guarantees that every uploaded file contributes chunks, regardless of global similarity rank."""
    out: list[Document] = []
    seen: set[str] = set()
    for src in sources:
        try:
            result = raw_coll.query(
                query_embeddings=[query_embedding],
                n_results=per_src_k,
                where={"source": src},
                include=["documents", "metadatas"],
            )
            docs_lists = result.get("documents") or [[]]
            meta_lists = result.get("metadatas") or [[]]
            for content, meta in zip(docs_lists[0], meta_lists[0]):
                if not content or not content.strip():
                    continue
                key = content[:80]
                if key not in seen:
                    seen.add(key)
                    out.append(Document(page_content=content, metadata=meta or {}))
        except Exception:
            pass
    return out


@app.get("/sources/{user_id}")
def list_sources(user_id: str) -> dict[str, Any]:
    """Return the list of source filenames currently indexed in Chroma for this user."""
    name = collection_name(user_id)
    try:
        client = chromadb.PersistentClient(path=CHROMA_PATH)
        names = {c.name for c in client.list_collections()}
        if name not in names:
            return {"sources": [], "chunk_count": 0}
        raw_coll = client.get_collection(name)
        all_meta = raw_coll.get(include=["metadatas"], limit=5000)
        sources = sorted({
            m.get("source", "") for m in (all_meta.get("metadatas") or []) if m.get("source")
        })
        chunk_count = len(all_meta.get("ids") or [])
        return {"sources": sources, "chunk_count": chunk_count}
    except Exception as exc:
        return {"sources": [], "chunk_count": 0, "error": str(exc)}


@app.post("/search")
def search(body: SearchRequest) -> dict[str, Any]:
    if VECTOR_BACKEND != "chroma":
        raise HTTPException(400, f"Unsupported VECTOR_BACKEND={VECTOR_BACKEND!r}.")

    name = collection_name(body.user_id)
    client = chromadb.PersistentClient(path=CHROMA_PATH)
    names = {c.name for c in client.list_collections()}
    if name not in names:
        return {"context": "", "image_refs": []}

    store = Chroma(
        persist_directory=CHROMA_PATH,
        embedding_function=embeddings,
        collection_name=name,
    )
    raw_coll = client.get_collection(name)

    # Discover all unique source documents in this user's collection.
    try:
        all_meta = raw_coll.get(include=["metadatas"], limit=5000)
        sources = list({m.get("source", "") for m in (all_meta.get("metadatas") or []) if m.get("source")})
    except Exception:
        sources = []

    # Embed query once, then search every source document independently.
    # This prevents any single document from monopolising the result set.
    try:
        query_vec = embeddings.embed_query(body.query)
        per_src_k = max(2, body.k // max(len(sources), 1))
        hits = _multi_source_hits(raw_coll, query_vec, sources, per_src_k)
    except Exception:
        hits = []

    # Fallback: plain similarity search when per-source retrieval fails.
    if not hits:
        hits = store.similarity_search(body.query, k=body.k)
    parts: list[str] = []
    image_refs: list[dict[str, Any]] = []
    seen_images: set[str] = set()
    uid = re.sub(r"[^a-zA-Z0-9_@.-]", "_", body.user_id.strip().lower())[:120]

    for d in hits:
        src = d.metadata.get("source", "?")
        mod = d.metadata.get("module", "")
        sec = str(d.metadata.get("section", "") or "").replace("\n", " ")
        pg = d.metadata.get("page")
        head = f"Source: {src}"
        if mod:
            head += f" | module: {mod}"
        if pg is not None:
            head += f" | page: {pg}"
        if sec:
            head += f" | section: {sec[:160]}"
        parts.append(f"{head}\n{d.page_content}")

        # Attach stored slide images for this chunk's page number
        if pg is not None and src != "?":
            safe = _safe_doc_name(src)
            img_dir = os.path.join(IMAGES_PATH, uid, safe)
            if os.path.isdir(img_dir):
                for img_file in sorted(os.listdir(img_dir)):
                    if img_file.startswith(f"slide_{pg}_"):
                        key = f"{src}:{pg}:{img_file}"
                        if key not in seen_images:
                            seen_images.add(key)
                            image_refs.append({
                                "source": src,
                                "slide": pg,
                                "url": f"/images/{uid}/{safe}/{img_file}",
                            })

    return {"context": "\n\n══════════════\n\n".join(parts), "image_refs": image_refs}
